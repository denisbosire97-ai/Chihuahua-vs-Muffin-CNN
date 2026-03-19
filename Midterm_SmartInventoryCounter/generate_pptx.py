from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

slides_data = [
    {
        "title": "Denis Aosa Safetycoordinator Glasses",
        "content": "Team Member: Denis Aosa\nCourse: ITAI 1378 - Computer Vision and AI\nTier: Midterm Project",
        "layout": 0
    },
    {
        "title": "The Problem",
        "content": "What real-world problem are you solving?\nDelivery drivers (like Amazon DSPs) face unpredictably dangerous 'Last 100 Yards' conditions. Trips on broken stairs, loose pets, and poor back mechanics cost logistics companies millions.\n\nWhy is it important?\nAccidents and injuries reduce driver retention and delay crucial deliveries.",
        "layout": 1
    },
    {
        "title": "Your Solution (Overview)",
        "content": "What will your system do?\nThe Denis Aosa Safetycoordinator Glasses are a wearable CV system that evaluates real-time camera feeds to alert drivers of high-risk situations.\n\nHow will it solve the problem?\nIt identifies critical hazards (stairs, pets) and poorly lit conditions, triggering a 'Stop & Assess' HUD visual. It also tracks spinal posture to ensure safe package lifting mechanics.",
        "layout": 1
    },
    {
        "title": "Technical Approach",
        "content": "Computer Vision Architectures:\n\n1. Object Detection: YOLOv8-Nano\nExtremely lightweight (under 4MB) to easily run on edge wearable constraints while detecting 5 key hazard classes.\n\n2. Pose Estimation: MediaPipe & YOLO-Pose\nMaps driver's skeletal angle to warn them if they bend at the waist instead of the knees.",
        "layout": 1
    },
    {
        "title": "System Diagram",
        "content": "Workflow:\n\n1. Input: Driver's POV camera feed.\n2. Perception: YOLOv8-Nano detects bounding boxes; YOLO-Pose maps skeletal nodes.\n3. Logic: safety_logic.py cross-references detected items with lighting conditions.\n4. Output (HUD): Displays visual cues (Red Box = STOP, Green Check = Safe Posture).\n5. Reporting: Generates an automatic Incident Log for the Safety Coordinator.",
        "layout": 1
    },
    {
        "title": "Success Metrics",
        "content": "Primary Metrics:\n- Recall: > 95% for Critical Hazards (False Positives > False Negatives)\n- Pose Accuracy: mAP > 85%\n\nSecondary Metrics:\n- Telemetry: Inference speed < 30ms (Real-time capability on edge processors)",
        "layout": 1
    },
    {
        "title": "Resources Needed",
        "content": "Compute:\n- Edge TPUs, Google Colab for initial training.\n\nFrameworks & Languages:\n- Python 3.10+\n- PyTorch, Ultralytics YOLOv8 logic\n- OpenCV\n\nEstimated Cost: Low (Software prototype)",
        "layout": 1
    }
]

for slide_data in slides_data:
    slide_layout = prs.slide_layouts[slide_data["layout"]]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = slide_data["title"]
    
    if slide_data["layout"] == 0:
        subtitle = slide.placeholders[1]
        subtitle.text = slide_data["content"]
    else:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        
        paragraphs = slide_data["content"].split('\n')
        tf.text = paragraphs[0]
        for p in paragraphs[1:]:
            p_obj = tf.add_paragraph()
            p_obj.text = p
            p_obj.font.size = Pt(20)

prs.save("C:/Users/denis.aosa/.gemini/antigravity/playground/vector-perihelion/Presentation_DenisAosa_Safetycoordinator.pptx")
print("SAFETYCOORDINATOR GLASSES PPTX GENERATED SUCCESSFULLY")
